from __future__ import annotations

import io
import hashlib
from dataclasses import dataclass
from typing import Any, Dict, Iterable, List, Sequence, Tuple

import numpy as np
from openai import OpenAI
from pypdf import PdfReader
from pptx import Presentation

from blooket_generator import ensure_api_key

MAX_FILE_SIZE_MB = 8
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100
MAX_CHUNKS_PER_FILE = 200
EMBEDDING_MODEL = "text-embedding-3-small"


@dataclass
class DocumentChunk:
    content: str
    source: str
    order: int
    metadata: Dict[str, Any]


def _normalize_text(value: str) -> str:
    return " ".join(value.split())


def _split_into_chunks(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    normalized = _normalize_text(text)
    if not normalized:
        return []

    chunks: List[str] = []
    start = 0
    while start < len(normalized):
        end = min(len(normalized), start + chunk_size)
        chunk = normalized[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == len(normalized):
            break
        start = max(0, end - overlap)
        if start >= len(normalized):
            break
    return chunks


def _extract_pdf_segments(data: bytes, file_name: str) -> List[Tuple[str, Dict[str, Any]]]:
    reader = PdfReader(io.BytesIO(data))
    segments: List[Tuple[str, Dict[str, Any]]] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if not text:
            continue
        segments.append((text, {"location": f"{file_name} p.{idx}"}))
    return segments


def _extract_pptx_segments(data: bytes, file_name: str) -> List[Tuple[str, Dict[str, Any]]]:
    presentation = Presentation(io.BytesIO(data))
    segments: List[Tuple[str, Dict[str, Any]]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        collected: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted = shape.text or ""
                if extracted:
                    collected.append(extracted)
        slide_text = "\n".join(collected).strip()
        if not slide_text:
            continue
        segments.append((slide_text, {"location": f"{file_name} slide {idx}"}))
    return segments


def _extract_segments(data: bytes, file_name: str) -> List[Tuple[str, Dict[str, Any]]]:
    lower_name = file_name.lower()
    if lower_name.endswith(".pdf"):
        return _extract_pdf_segments(data, file_name)
    if lower_name.endswith(".pptx"):
        return _extract_pptx_segments(data, file_name)
    raise ValueError("지원하지 않는 파일 형식입니다. PDF 또는 PPTX 파일만 업로드하세요.")


def _build_document_chunks(
    segments: Sequence[Tuple[str, Dict[str, Any]]],
    file_name: str,
    file_hash: str,
) -> List[DocumentChunk]:
    chunks: List[DocumentChunk] = []
    order = 0
    for text, metadata in segments:
        for chunk_text in _split_into_chunks(text):
            order += 1
            enriched_meta = dict(metadata)
            enriched_meta.update({"file_hash": file_hash, "source": file_name})
            chunks.append(
                DocumentChunk(
                    content=chunk_text,
                    source=file_name,
                    order=order,
                    metadata=enriched_meta,
                )
            )
    return chunks


def _embed_texts(texts: Sequence[str], api_key: str, batch_size: int = 64) -> np.ndarray:
    client = OpenAI(api_key=ensure_api_key(api_key))
    vectors: List[List[float]] = []
    for start in range(0, len(texts), batch_size):
        batch = texts[start : start + batch_size]
        response = client.embeddings.create(model=EMBEDDING_MODEL, input=list(batch))
        ordered = sorted(response.data, key=lambda item: item.index)
        vectors.extend(item.embedding for item in ordered)
    return np.asarray(vectors, dtype=np.float32)


def process_uploaded_documents(
    uploads: Sequence[Any],
    api_key: str,
    index: Dict[str, Any] | None = None,
) -> Tuple[Dict[str, Any], Dict[str, List[str]]]:
    index = index or {
        "vectors": None,
        "chunks": [],
        "hashes": set(),
        "sources": {},
        "model": EMBEDDING_MODEL,
    }

    status = {"added": [], "warnings": [], "errors": []}

    for uploaded in uploads:
        if not hasattr(uploaded, "name") or not hasattr(uploaded, "getvalue"):
            continue
        name = uploaded.name
        data = uploaded.getvalue()
        file_hash = hashlib.md5(data).hexdigest()

        if file_hash in index["hashes"]:
            continue

        if len(data) > MAX_FILE_SIZE_BYTES:
            status["errors"].append(
                f"{name}: 파일 크기가 {MAX_FILE_SIZE_MB}MB 제한을 초과합니다."
            )
            continue

        try:
            segments = _extract_segments(data, name)
        except Exception as exc:  # noqa: BLE001
            status["errors"].append(f"{name}: 텍스트 추출에 실패했습니다 ({exc}).")
            continue

        if not segments:
            status["warnings"].append(f"{name}: 추출할 텍스트가 없어 건너뜁니다.")
            continue

        chunks = _build_document_chunks(segments, name, file_hash)
        if not chunks:
            status["warnings"].append(f"{name}: 생성된 청크가 없어 건너뜁니다.")
            continue

        truncated = False
        if len(chunks) > MAX_CHUNKS_PER_FILE:
            chunks = chunks[:MAX_CHUNKS_PER_FILE]
            truncated = True

        try:
            vectors = _embed_texts([chunk.content for chunk in chunks], api_key)
        except Exception as exc:  # noqa: BLE001
            status["errors"].append(f"{name}: 임베딩 생성 중 오류가 발생했습니다 ({exc}).")
            continue

        if index["vectors"] is None:
            index["vectors"] = vectors
        else:
            index["vectors"] = np.vstack([index["vectors"], vectors])

        index["chunks"].extend(chunks)
        index["hashes"].add(file_hash)
        index["sources"][file_hash] = {
            "name": name,
            "segments": len(segments),
            "chunks": len(chunks),
            "truncated": truncated,
        }

        message = f"{name}: {len(chunks)}개 청크"
        if truncated:
            message += " (일부만 사용)"
        status["added"].append(message)

    return index, status


def retrieve_relevant_context(
    query: str,
    index: Dict[str, Any],
    api_key: str,
    top_k: int = 5,
    min_score: float = 0.2,
) -> str:
    if not query.strip():
        return ""
    if not index or not index.get("chunks") or index.get("vectors") is None:
        return ""

    client = OpenAI(api_key=ensure_api_key(api_key))
    response = client.embeddings.create(model=EMBEDDING_MODEL, input=[query])
    query_vector = np.asarray(response.data[0].embedding, dtype=np.float32)

    vectors: np.ndarray = index["vectors"]
    similarities = vectors @ query_vector
    query_norm = float(np.linalg.norm(query_vector)) or 1.0
    vector_norms = np.linalg.norm(vectors, axis=1)
    similarities = similarities / (vector_norms * query_norm + 1e-8)

    ranked = np.argsort(similarities)[::-1]
    selected_parts: List[str] = []
    used: set[Tuple[str, int]] = set()

    for idx in ranked[:top_k]:
        score = float(similarities[idx])
        if score < min_score:
            continue
        chunk: DocumentChunk = index["chunks"][idx]
        key = (chunk.metadata.get("source", chunk.source), chunk.order)
        if key in used:
            continue
        used.add(key)
        location = chunk.metadata.get("location", chunk.source)
        selected_parts.append(f"[{location}] {chunk.content}")

    return "\n\n".join(selected_parts)


__all__ = [
    "DocumentChunk",
    "EMBEDDING_MODEL",
    "MAX_FILE_SIZE_MB",
    "process_uploaded_documents",
    "retrieve_relevant_context",
]
